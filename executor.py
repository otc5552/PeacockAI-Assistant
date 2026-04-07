import os
import subprocess
import base64
import json
import shutil
import time
import io
import urllib.request
import urllib.parse
from pathlib import Path
from typing import Optional

# ── Desktop path helper ──────────────────────────────────────────────────────
def get_desktop():
    return str(Path.home() / "Desktop")


# ══════════════════════════════════════════════════════════════════════════════
#  APP CONTROL
# ══════════════════════════════════════════════════════════════════════════════
def open_app(app_name: str) -> str:
    """Open an application by name or path."""
    known = {
        "vs code": "code",
        "vscode": "code",
        "visual studio code": "code",
        "android studio": "studio64.exe",
        "chrome": "chrome",
        "google chrome": "chrome",
        "firefox": "firefox",
        "notepad": "notepad",
        "notepad++": "notepad++",
        "photoshop": "photoshop",
        "word": "winword",
        "excel": "excel",
        "powerpoint": "powerpnt",
        "explorer": "explorer",
        "cmd": "cmd",
        "powershell": "powershell",
        "calculator": "calc",
        "paint": "mspaint",
        "task manager": "taskmgr",
    }
    cmd = known.get(app_name.lower().strip(), app_name)
    try:
        subprocess.Popen(cmd, shell=True)
        return f"✅ Opening {app_name}..."
    except Exception as e:
        return f"❌ Could not open {app_name}: {e}"


# ══════════════════════════════════════════════════════════════════════════════
#  SCREEN & INPUT CONTROL
# ══════════════════════════════════════════════════════════════════════════════
def take_screenshot() -> dict:
    """Take a screenshot and return base64."""
    try:
        import pyautogui
        from PIL import Image
        screenshot = pyautogui.screenshot()
        buf = io.BytesIO()
        screenshot.save(buf, format="PNG")
        b64 = base64.b64encode(buf.getvalue()).decode()
        return {"success": True, "image": b64, "message": "Screenshot taken"}
    except Exception as e:
        return {"success": False, "message": f"Screenshot failed: {e}"}


def mouse_click(x: int, y: int, button: str = "left") -> str:
    try:
        import pyautogui
        if button == "double":
            pyautogui.doubleClick(x, y)
        elif button == "right":
            pyautogui.rightClick(x, y)
        else:
            pyautogui.click(x, y)
        return f"✅ Clicked at ({x}, {y})"
    except Exception as e:
        return f"❌ Click failed: {e}"


def keyboard_type(text: str, press_enter: bool = False) -> str:
    try:
        import pyautogui
        pyautogui.write(text, interval=0.02)
        if press_enter:
            pyautogui.press("enter")
        return f"✅ Typed: {text}"
    except Exception as e:
        return f"❌ Type failed: {e}"


def keyboard_shortcut(keys: str) -> str:
    try:
        import pyautogui
        key_list = [k.strip() for k in keys.lower().split("+")]
        pyautogui.hotkey(*key_list)
        return f"✅ Pressed: {keys}"
    except Exception as e:
        return f"❌ Shortcut failed: {e}"


def scroll(amount: int, x: int = None, y: int = None) -> str:
    try:
        import pyautogui
        if x and y:
            pyautogui.scroll(amount, x=x, y=y)
        else:
            pyautogui.scroll(amount)
        return f"✅ Scrolled {amount}"
    except Exception as e:
        return f"❌ Scroll failed: {e}"


# ══════════════════════════════════════════════════════════════════════════════
#  FILE SYSTEM
# ══════════════════════════════════════════════════════════════════════════════
def read_file(path: str) -> str:
    try:
        with open(path, "r", encoding="utf-8", errors="replace") as f:
            content = f.read()
        return f"✅ File content:\n{content[:5000]}" + ("\n[...truncated]" if len(content) > 5000 else "")
    except Exception as e:
        return f"❌ Could not read file: {e}"


def write_file(path: str, content: str) -> str:
    try:
        os.makedirs(os.path.dirname(os.path.abspath(path)), exist_ok=True)
        with open(path, "w", encoding="utf-8") as f:
            f.write(content)
        return f"✅ File saved: {path}"
    except Exception as e:
        return f"❌ Could not write file: {e}"


def list_directory(path: str = None) -> str:
    try:
        target = path or get_desktop()
        items = os.listdir(target)
        result = f"📁 Contents of {target}:\n"
        for item in sorted(items):
            full = os.path.join(target, item)
            icon = "📁" if os.path.isdir(full) else "📄"
            result += f"  {icon} {item}\n"
        return result
    except Exception as e:
        return f"❌ Could not list directory: {e}"


def delete_file(path: str) -> str:
    try:
        if os.path.isdir(path):
            shutil.rmtree(path)
        else:
            os.remove(path)
        return f"✅ Deleted: {path}"
    except Exception as e:
        return f"❌ Could not delete: {e}"


# ══════════════════════════════════════════════════════════════════════════════
#  WEB SEARCH
# ══════════════════════════════════════════════════════════════════════════════
def search_web(query: str) -> str:
    try:
        encoded = urllib.parse.quote(query)
        url = f"https://html.duckduckgo.com/html/?q={encoded}"
        req = urllib.request.Request(url, headers={"User-Agent": "Mozilla/5.0"})
        with urllib.request.urlopen(req, timeout=10) as resp:
            html = resp.read().decode("utf-8", errors="replace")

        import re
        results = re.findall(r'class="result__snippet"[^>]*>(.*?)</a>', html, re.DOTALL)
        titles = re.findall(r'class="result__a"[^>]*>(.*?)</a>', html, re.DOTALL)
        clean = lambda s: re.sub(r'<[^>]+>', '', s).strip()

        output = f"🔍 Search results for: {query}\n\n"
        for i, (t, r) in enumerate(zip(titles[:5], results[:5]), 1):
            output += f"{i}. **{clean(t)}**\n{clean(r)}\n\n"
        return output if titles else f"No results found for: {query}"
    except Exception as e:
        return f"❌ Search failed: {e}"


# ══════════════════════════════════════════════════════════════════════════════
#  OFFICE DOCUMENTS
# ══════════════════════════════════════════════════════════════════════════════
def create_word_document(filename: str, content: str, title: str = "", save_path: str = None) -> str:
    try:
        from docx import Document
        from docx.shared import Pt, RGBColor
        from docx.enum.text import WD_ALIGN_PARAGRAPH

        doc = Document()

        if title:
            h = doc.add_heading(title, 0)
            h.alignment = WD_ALIGN_PARAGRAPH.CENTER

        for para in content.split("\n"):
            if para.strip().startswith("# "):
                doc.add_heading(para[2:], level=1)
            elif para.strip().startswith("## "):
                doc.add_heading(para[3:], level=2)
            elif para.strip().startswith("- ") or para.strip().startswith("* "):
                doc.add_paragraph(para.strip()[2:], style="List Bullet")
            elif para.strip():
                doc.add_paragraph(para)

        out_dir = save_path or get_desktop()
        os.makedirs(out_dir, exist_ok=True)
        if not filename.endswith(".docx"):
            filename += ".docx"
        out_path = os.path.join(out_dir, filename)
        doc.save(out_path)
        return f"✅ Word document created: {out_path}"
    except ImportError:
        return "❌ python-docx not installed. Run: pip install python-docx"
    except Exception as e:
        return f"❌ Word creation failed: {e}"


def create_powerpoint(filename: str, slides: list, title: str = "", save_path: str = None) -> str:
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt
        from pptx.dml.color import RGBColor

        prs = Presentation()
        prs.slide_width = Inches(13.33)
        prs.slide_height = Inches(7.5)

        # Title slide
        if title:
            slide_layout = prs.slide_layouts[0]
            slide = prs.slides.add_slide(slide_layout)
            slide.shapes.title.text = title
            if slide.placeholders[1]:
                slide.placeholders[1].text = "PeacockAgent"

        # Content slides
        for s in slides:
            layout = prs.slide_layouts[1]
            slide = prs.slides.add_slide(layout)
            slide.shapes.title.text = s.get("title", "")
            body = slide.placeholders[1]
            tf = body.text_frame
            tf.text = s.get("content", "")

        out_dir = save_path or get_desktop()
        os.makedirs(out_dir, exist_ok=True)
        if not filename.endswith(".pptx"):
            filename += ".pptx"
        out_path = os.path.join(out_dir, filename)
        prs.save(out_path)
        return f"✅ PowerPoint created: {out_path}"
    except ImportError:
        return "❌ python-pptx not installed. Run: pip install python-pptx"
    except Exception as e:
        return f"❌ PowerPoint creation failed: {e}"


def create_pdf(filename: str, content: str, title: str = "", save_path: str = None) -> str:
    try:
        from reportlab.lib.pagesizes import A4
        from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
        from reportlab.lib.units import cm
        from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer
        from reportlab.lib import colors

        out_dir = save_path or get_desktop()
        os.makedirs(out_dir, exist_ok=True)
        if not filename.endswith(".pdf"):
            filename += ".pdf"
        out_path = os.path.join(out_dir, filename)

        doc = SimpleDocTemplate(out_path, pagesize=A4,
                                leftMargin=2*cm, rightMargin=2*cm,
                                topMargin=2*cm, bottomMargin=2*cm)
        styles = getSampleStyleSheet()
        story = []

        if title:
            story.append(Paragraph(title, styles["Title"]))
            story.append(Spacer(1, 0.5*cm))

        for line in content.split("\n"):
            if line.strip():
                story.append(Paragraph(line, styles["Normal"]))
                story.append(Spacer(1, 0.2*cm))

        doc.build(story)
        return f"✅ PDF created: {out_path}"
    except ImportError:
        return "❌ reportlab not installed. Run: pip install reportlab"
    except Exception as e:
        return f"❌ PDF creation failed: {e}"


# ══════════════════════════════════════════════════════════════════════════════
#  IMAGE GENERATION & EDITING
# ══════════════════════════════════════════════════════════════════════════════
def generate_image(prompt: str, width: int = 1024, height: int = 1024, save_path: str = None) -> str:
    try:
        encoded = urllib.parse.quote(prompt)
        url = f"https://image.pollinations.ai/prompt/{encoded}?width={width}&height={height}&nologo=true"

        out_dir = save_path or get_desktop()
        os.makedirs(out_dir, exist_ok=True)
        safe_name = "".join(c for c in prompt[:30] if c.isalnum() or c in " _-").strip().replace(" ", "_")
        filename = f"{safe_name}_{int(time.time())}.png"
        out_path = os.path.join(out_dir, filename)

        req = urllib.request.Request(url, headers={"User-Agent": "Mozilla/5.0"})
        with urllib.request.urlopen(req, timeout=60) as resp:
            with open(out_path, "wb") as f:
                f.write(resp.read())

        return f"✅ Image generated and saved: {out_path}"
    except Exception as e:
        return f"❌ Image generation failed: {e}"


def edit_image(input_path: str, operations: list, output_path: str = None) -> str:
    try:
        from PIL import Image, ImageEnhance, ImageFilter

        img = Image.open(input_path)

        for op in operations:
            op_type = op.get("type", "")
            params = op.get("params", {})

            if op_type == "resize":
                w = params.get("width", img.width)
                h = params.get("height", img.height)
                img = img.resize((int(w), int(h)), Image.LANCZOS)

            elif op_type == "crop":
                left = params.get("left", 0)
                top = params.get("top", 0)
                right = params.get("right", img.width)
                bottom = params.get("bottom", img.height)
                img = img.crop((int(left), int(top), int(right), int(bottom)))

            elif op_type == "rotate":
                angle = params.get("angle", 90)
                img = img.rotate(int(angle), expand=True)

            elif op_type == "brightness":
                factor = params.get("factor", 1.2)
                img = ImageEnhance.Brightness(img).enhance(float(factor))

            elif op_type == "contrast":
                factor = params.get("factor", 1.2)
                img = ImageEnhance.Contrast(img).enhance(float(factor))

            elif op_type == "grayscale":
                img = img.convert("L").convert("RGB")

            elif op_type == "blur":
                radius = params.get("radius", 2)
                img = img.filter(ImageFilter.GaussianBlur(radius=int(radius)))

            elif op_type == "sharpen":
                img = img.filter(ImageFilter.SHARPEN)

            elif op_type == "flip":
                direction = params.get("direction", "horizontal")
                if direction == "horizontal":
                    img = img.transpose(Image.FLIP_LEFT_RIGHT)
                else:
                    img = img.transpose(Image.FLIP_TOP_BOTTOM)

        out = output_path or input_path.replace(".", "_edited.")
        img.save(out)
        return f"✅ Image edited and saved: {out}"
    except ImportError:
        return "❌ Pillow not installed. Run: pip install Pillow"
    except Exception as e:
        return f"❌ Image editing failed: {e}"


# ══════════════════════════════════════════════════════════════════════════════
#  TERMINAL COMMAND
# ══════════════════════════════════════════════════════════════════════════════
def run_command(command: str, shell: str = "cmd") -> str:
    try:
        if shell == "powershell":
            result = subprocess.run(
                ["powershell", "-Command", command],
                capture_output=True, text=True, timeout=30
            )
        else:
            result = subprocess.run(
                command, shell=True, capture_output=True, text=True, timeout=30
            )
        output = result.stdout.strip() or result.stderr.strip()
        return f"✅ Command output:\n{output[:3000]}" if output else "✅ Command executed (no output)"
    except subprocess.TimeoutExpired:
        return "⚠️ Command timed out after 30 seconds"
    except Exception as e:
        return f"❌ Command failed: {e}"


# ══════════════════════════════════════════════════════════════════════════════
#  TOOL DISPATCHER
# ══════════════════════════════════════════════════════════════════════════════
def dispatch_tool(tool_name: str, tool_args: dict) -> str:
    tools_map = {
        "open_app": lambda a: open_app(a["app_name"]),
        "take_screenshot": lambda a: take_screenshot(),
        "mouse_click": lambda a: mouse_click(a["x"], a["y"], a.get("button", "left")),
        "keyboard_type": lambda a: keyboard_type(a["text"], a.get("press_enter", False)),
        "keyboard_shortcut": lambda a: keyboard_shortcut(a["keys"]),
        "scroll": lambda a: scroll(a["amount"], a.get("x"), a.get("y")),
        "read_file": lambda a: read_file(a["path"]),
        "write_file": lambda a: write_file(a["path"], a["content"]),
        "list_directory": lambda a: list_directory(a.get("path")),
        "delete_file": lambda a: delete_file(a["path"]),
        "search_web": lambda a: search_web(a["query"]),
        "create_word_document": lambda a: create_word_document(
            a["filename"], a["content"], a.get("title", ""), a.get("save_path")),
        "create_powerpoint": lambda a: create_powerpoint(
            a["filename"], a["slides"], a.get("title", ""), a.get("save_path")),
        "create_pdf": lambda a: create_pdf(
            a["filename"], a["content"], a.get("title", ""), a.get("save_path")),
        "generate_image": lambda a: generate_image(
            a["prompt"], a.get("width", 1024), a.get("height", 1024), a.get("save_path")),
        "edit_image": lambda a: edit_image(
            a["input_path"], a["operations"], a.get("output_path")),
        "run_command": lambda a: run_command(a["command"], a.get("shell", "cmd")),
    }

    fn = tools_map.get(tool_name)
    if fn:
        return fn(tool_args)
    return f"❌ Unknown tool: {tool_name}"
